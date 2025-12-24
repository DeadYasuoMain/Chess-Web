from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from datetime import datetime

prs = Presentation()

slides = [
    {
        "title": "HTML, CSS & JavaScript",
        "subtitle": "Industrial Training — Introduction to Web Development",
        "notes": "Introduce yourself, goals: explain roles of HTML/CSS/JS and build a small demo by the end of session."
    },
    {
        "title": "Training Objectives",
        "bullets": [
            "Understand structure (HTML)",
            "Style and layout with CSS",
            "Behavior and interactivity with JavaScript",
            "Hands-on mini project and resources"
        ],
        "notes": "Explain the expected takeaways and the format (theory + demos + exercise)."
    },
    {
        "title": "How the Web Works",
        "bullets": ["Browser, Server, HTTP","HTML = structure","CSS = presentation","JS = behavior"],
        "notes": "Give a short picture of requests, responses and how browsers render pages."
    },
    {
        "title": "Introduction to HTML",
        "bullets": ["Tags and elements","Document structure: doctype, html, head, body","Attributes and semantic tags"],
        "code": "<!DOCTYPE html>\n<html>\n  <head>\n    <title>My Page</title>\n  </head>\n  <body>\n    <h1>Hello World</h1>\n  </body>\n</html>",
        "notes": "Show a simple HTML page. Explain tree structure and semantic importance."
    },
    {
        "title": "Common HTML Elements",
        "bullets": ["Headings, paragraphs, links, images","Lists and tables","Forms and inputs"],
        "notes": "Show quick examples and talk about when to use which element."
    },
    {
        "title": "HTML Best Practices & Accessibility",
        "bullets": ["Use semantic tags","Provide alt text for images","Proper heading order & keyboard accessibility"],
        "notes": "Emphasize accessibility and SEO benefits."
    },
    {
        "title": "Introduction to CSS",
        "bullets": ["Selectors & properties","Cascade & specificity","Box model (margin/border/padding/content)"],
        "code": "h1 {\n  color: #2b8a3e;\n  font-family: Arial, sans-serif;\n}",
        "notes": "Explain how styles are attached (inline, internal, external) and specificity."
    },
    {
        "title": "Layout Techniques",
        "bullets": ["Float (legacy)","Flexbox: one-dimensional layouts","Grid: two-dimensional layouts","Responsive: media queries"],
        "notes": "Show when to use Flexbox vs Grid and simple responsive tips."
    },
    {
        "title": "Styling Examples",
        "bullets": ["Typography and spacing","Colors and themes","Hover/focus states and transitions"],
        "notes": "Demo a small style change live in the workshop."
    },
    {
        "title": "Introduction to JavaScript",
        "bullets": ["Variables, functions, events","DOM: query and manipulate nodes","Use JS to respond to user actions"],
        "code": "document.getElementById('btn').addEventListener('click', () => {\n  alert('Clicked!');\n});",
        "notes": "Explain JS role and safety (avoid blocking main thread)."
    },
    {
        "title": "DOM & Events",
        "bullets": ["Selecting elements","Modifying content and attributes","Handling events like click/submit"],
        "notes": "Show sample: toggle a theme or simple form validation."
    },
    {
        "title": "Modern JS & Toolchain",
        "bullets": ["ES6 features: let/const, arrow functions","Dev tools and debugging","Brief mention: bundlers, npm, frameworks"],
        "notes": "Keep it high-level and point to resources for deeper learning."
    },
    {
        "title": "Mini Project Idea",
        "bullets": ["Build a theme toggle or todo list","HTML for structure, CSS for style, JS for behavior","Exercise + stretch goals"],
        "notes": "Explain expected time and how to split tasks for attendees."
    },
    {
        "title": "Testing, Debugging & Deployment",
        "bullets": ["Use browser dev tools","Linting (HTML/CSS/JS)","Deploy: GitHub Pages, Netlify"],
        "notes": "Show a quick demo of opening dev tools and inspecting elements."
    },
    {
        "title": "Resources & Q&A",
        "bullets": ["MDN, W3Schools, freeCodeCamp","Cheat-sheets and exercise repo","Questions and next steps"],
        "notes": "Provide links and invite questions."
    }
]

# Helper to add a title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = slides[0]["title"]
try:
    slide.placeholders[1].text = slides[0]["subtitle"] + "\n\n" + datetime.today().strftime('%B %d, %Y')
except Exception:
    pass
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = slides[0]["notes"]

# Add the rest of slides
for s in slides[1:]:
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = s["title"]

    # Add bullets if present
    if "bullets" in s:
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(s["bullets"]):
            if i == 0:
                p = body.paragraphs[0]
                p.text = b
            else:
                p = body.add_paragraph()
                p.text = b
            p.level = 0

    # Add code box if present
    if "code" in s:
        left = Inches(1)
        top = Inches(3.0)
        width = Inches(8)
        height = Inches(1.8)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.name = 'Consolas'
        p.font.size = Pt(12)
        p.text = s["code"]

    # Speaker notes
    if "notes" in s:
        ns = slide.notes_slide.notes_text_frame
        ns.text = s["notes"]

# Final resources slide: add contact info
final_slide = prs.slides.add_slide(prs.slide_layouts[1])
final_slide.shapes.title.text = "Contact & Next Steps"
tf = final_slide.shapes.placeholders[1].text_frame
tf.clear()
para = tf.paragraphs[0]
para.text = "Instructor: Your Name"
para.level = 0
p2 = tf.add_paragraph()
p2.text = "Repo: (add link) | Slides: available after session"
p2.level = 0

output_path = 'HTML_CSS_JS_Industrial_Training.pptx'
prs.save(output_path)
print(f"Saved presentation to {output_path}")
